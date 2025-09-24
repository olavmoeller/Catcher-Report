import matplotlib.pyplot as plt
import matplotlib.patheffects as path_effects
from matplotlib.patches import Rectangle
from matplotlib import colors
import pandas as pd
import numpy as np
from scipy.stats import gaussian_kde
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re
from tkinter import filedialog
import comtypes.client
import os

# Team toggle - set to 'beavers' or 'knights'
team = 'knights'  # Change this to switch teams

# Insert date (or what you want the top right to show)
date = 'Beavers Season'

# Insert sb/attempt data
sb_att_dict = {
    # "Catcher Name": (Steals, Attempts),
    "Rios, Xavier": (23, 34),
    "Fernandez, Ian": (10, 17),
    "Hazen, Aidan": (1, 1),
    "Marquez, Jacob": (4, 7),
    "Weber, Wilson": (50, 68),
    "Hubbard, Bryce": (4, 5)
    # Add more catchers as needed
}

if team == 'beavers':
    team_code = 'ORE_BEA'
    template_file = "Catching_Template.pptx"
else:  # knights
    team_code = 'COR_KNI'
    template_file = "Catching_Template_Knights.pptx"

# Strike zone constants (matching umpire accuracy report)
STRIKE_ZONE_WIDTH_IN = 17  # inches
STRIKE_ZONE_TOP_IN = 39    # inches
STRIKE_ZONE_BOT_IN = 19.7  # inches
STRIKE_ZONE_LEFT_IN = -STRIKE_ZONE_WIDTH_IN / 2
STRIKE_ZONE_RIGHT_IN = STRIKE_ZONE_WIDTH_IN / 2

# Opens Window to Choose CSV files
csv_files = filedialog.askopenfilenames(title="Select CSV files", filetypes=(("CSV files", "*.csv"), ("all files", "*.*")))

# Check if files were selected
if csv_files:
    # Initialize an empty list to store DataFrames
    dfs = []

    # Iterate through selected CSV files and read them into DataFrames
    for file in csv_files:
        df = pd.read_csv(file)
        dfs.append(df)

    # Concatenate all DataFrames
    global_df = pd.concat(dfs, ignore_index=True)

    # Print the concatenated DataFrame or perform further operations
    print(global_df)
else:
    print("No CSV files selected.")

def processing_full(dataframe):
    # Filter called pitches
    df = dataframe[dataframe['PitchCall'].isin(['BallCalled', 'StrikeCalled'])].copy()

    if df.empty:
        return df  # Return early if there's nothing to process

    # Convert to inches and flip side for catcher's POV
    df['PlateLocHeightIn'] = 12 * df['PlateLocHeight']
    df['PlateLocSideIn'] = -12 * df['PlateLocSide']

    # Map calls to numeric values
    call_map = {'BallCalled': 0, 'StrikeCalled': 1}
    df['CallValue'] = df['PitchCall'].map(call_map)

    # Drop NAs for relevant columns
    df = df[['PitchCall', 'PlateLocHeightIn', 'PlateLocSideIn', 'Catcher', 'BatterSide', 'CallValue']].dropna()

    return df

def load_kde_cache(cache_file="kde_cache.npz"):
    """
    Load pre-computed KDE probability grids
    """
    if not os.path.exists(cache_file):
        print(f"❌ KDE cache file '{cache_file}' not found!")
        print("Please run 'create_kde_cache.py' first to create the cache from your master CSV.")
        return None
    
    print(f"Loading KDE cache from {cache_file}...")
    cache = np.load(cache_file)
    
    return {
        'league_strike_prob_l': cache['league_strike_prob_l'],
        'league_strike_prob_r': cache['league_strike_prob_r'],
        'x_grid': cache['x_grid'],
        'y_grid': cache['y_grid']
    }

def create_kde_heatmap_with_cache(catcher_data, league_cache, batter_side, ax):
    """
    Create KDE heatmap using pre-computed league data
    """
    if catcher_data.empty:
        # If no catcher data, show empty plot
        ax.set_xlim(-20, 20)
        ax.set_ylim(10, 50)
        ax.text(0, 30, f"No data for {batter_side} batters", ha='center', va='center', fontsize=14)
        return
    
    # Get pre-computed league data
    if batter_side == 'Left':
        league_strike_prob = league_cache['league_strike_prob_l']
    else:
        league_strike_prob = league_cache['league_strike_prob_r']
    
    x_grid = league_cache['x_grid']
    y_grid = league_cache['y_grid']
    X, Y = np.meshgrid(x_grid, y_grid)
    positions = np.vstack([X.ravel(), Y.ravel()])
    
    # Calculate catcher KDE
    catcher_strikes = catcher_data[catcher_data['CallValue'] == 1][['PlateLocSideIn', 'PlateLocHeightIn']].values.T
    catcher_balls = catcher_data[catcher_data['CallValue'] == 0][['PlateLocSideIn', 'PlateLocHeightIn']].values.T
    
    def calculate_strike_probability(strikes, balls, positions):
        if len(strikes) == 0 and len(balls) == 0:
            return np.zeros(positions.shape[1])
        
        total_pitches = len(strikes) + len(balls)
        if total_pitches == 0:
            return np.zeros(positions.shape[1])
        
        # KDE for strikes
        if len(strikes) > 0:
            kde_strikes = gaussian_kde(strikes, bw_method='silverman')
            strike_density = kde_strikes(positions)
        else:
            strike_density = np.zeros(positions.shape[1])
        
        # KDE for balls
        if len(balls) > 0:
            kde_balls = gaussian_kde(balls, bw_method='silverman')
            ball_density = kde_balls(positions)
        else:
            ball_density = np.zeros(positions.shape[1])
        
        # Calculate strike probability
        total_density = strike_density + ball_density
        strike_prob = np.where(total_density > 0, strike_density / total_density, 0)
        
        return strike_prob
    
    def calculate_confidence_map(catcher_data, positions, x_grid, y_grid):
        """Calculate confidence based on local sample density"""
        X, Y = np.meshgrid(x_grid, y_grid)
        confidence = np.zeros(positions.shape[1])
        
        for i, (x, y) in enumerate(zip(positions[0], positions[1])):
            # Count pitches within 2 inches of this point
            distances = np.sqrt((catcher_data['PlateLocSideIn'] - x)**2 + 
                              (catcher_data['PlateLocHeightIn'] - y)**2)
            nearby_pitches = np.sum(distances < 2.0)
            
            # Convert to confidence (0.3 to 1.0)
            confidence[i] = min(1.0, 0.3 + 0.7 * (nearby_pitches / 20))
        
        return confidence.reshape(X.shape)
    
    catcher_strike_prob = calculate_strike_probability(catcher_strikes, catcher_balls, positions)
    
    # Calculate difference from league average
    diff = catcher_strike_prob - league_strike_prob
    diff_grid = diff.reshape(X.shape)
    
    # Calculate confidence map
    confidence_map = calculate_confidence_map(catcher_data, positions, x_grid, y_grid)
    
    # Apply confidence to the difference values - low confidence fades toward zero (white)
    diff_grid_adjusted = diff_grid * confidence_map
    
    # Create heatmap
    cmap = colors.LinearSegmentedColormap.from_list('bwr', ["blue", "white", "red"], N=256)
    norm = colors.Normalize(-0.3, 0.3)  # Adjust range as needed
    
    im = ax.imshow(diff_grid_adjusted, extent=[-20, 20, 10, 50], origin='lower', 
                   cmap=cmap, norm=norm, aspect='auto', alpha=0.7)
    
    # Draw strike zone
    strike_rect = Rectangle((STRIKE_ZONE_LEFT_IN, STRIKE_ZONE_BOT_IN),
                           STRIKE_ZONE_WIDTH_IN, STRIKE_ZONE_TOP_IN - STRIKE_ZONE_BOT_IN,
                           linewidth=2, edgecolor='black', facecolor='none')
    ax.add_patch(strike_rect)
    
    # Draw shadow zone (1.5 inches around the strike zone - represents ball center for strikes)
    shadow_left = STRIKE_ZONE_LEFT_IN - 1.5
    shadow_right = STRIKE_ZONE_RIGHT_IN + 1.5
    shadow_bot = STRIKE_ZONE_BOT_IN - 1.5
    shadow_top = STRIKE_ZONE_TOP_IN + 1.5
    
    shadow_rect = Rectangle((shadow_left, shadow_bot),
                           shadow_right - shadow_left, shadow_top - shadow_bot,
                           linewidth=1, edgecolor='gray', facecolor='none', linestyle='--')
    ax.add_patch(shadow_rect)
    
    # Add batter position indicator
    x = 16 if batter_side == 'Left' else -16
    label = "LHB\nstands\nhere" if batter_side == 'Left' else "RHB\nstands\nhere"
    ax.text(x, 30, label, ha='center', va='center', color='black', fontsize=10,
            path_effects=[path_effects.Stroke(linewidth=2, foreground='white'), path_effects.Normal()])
    
    # Set plot limits and labels
    ax.set_xlim(-20, 20)
    ax.set_ylim(10, 50)
    ax.set_aspect('equal', adjustable='box')
    
    # Add colorbar
    cbar = plt.colorbar(im, ax=ax, shrink=0.8)
    cbar.set_label('Strike Rate vs League Average', rotation=270, labelpad=15)
    
    # Add title
    ax.set_title(f'Red = Better than League, Blue = Worse', 
                fontsize=12, pad=10)

def kde_plot_with_cache(name, df, league_cache, save_path=None):
    """
    Create KDE heatmap comparison for a catcher using cached league data
    """
    filt_df = df[(df['PitcherTeam'] == team_code) & (df['Catcher'] == name)]
    print(f"{name}: filtered dataframe shape = {filt_df.shape}")

    if filt_df.empty:
        print(f"No pitch data for catcher {name} after filtering.")
        return

    # Process data
    proc_df = processing_full(filt_df)
    
    if proc_df.empty:
        print(f"No processed data for catcher {name}.")
        return

    # Split by batter side
    df_l = proc_df[proc_df['BatterSide'] == 'Left']
    df_r = proc_df[proc_df['BatterSide'] == 'Right']

    # Check minimum sample sizes for KDE
    min_samples = 5  # Minimum pitches needed for meaningful KDE
    
    if len(df_l) < min_samples and len(df_r) < min_samples:
        print(f"Catcher {name} has insufficient data for KDE plots (LHB: {len(df_l)}, RHB: {len(df_r)}). Skipping.")
        return

    # Create subplots
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 6))
    
    # Create heatmaps using cache (only if enough data)
    if len(df_l) >= min_samples:
        create_kde_heatmap_with_cache(df_l, league_cache, 'Left', ax1)
    else:
        ax1.set_xlim(-20, 20)
        ax1.set_ylim(10, 50)
        ax1.text(0, 30, f"Insufficient LHB data\n({len(df_l)} pitches)", ha='center', va='center', fontsize=14)
        ax1.set_title('Left-Handed Batters', fontsize=12)
    
    if len(df_r) >= min_samples:
        create_kde_heatmap_with_cache(df_r, league_cache, 'Right', ax2)
    else:
        ax2.set_xlim(-20, 20)
        ax2.set_ylim(10, 50)
        ax2.text(0, 30, f"Insufficient RHB data\n({len(df_r)} pitches)", ha='center', va='center', fontsize=14)
        ax2.set_title('Right-Handed Batters', fontsize=12)

    plt.tight_layout()

    if save_path:
        fig.savefig(save_path, bbox_inches='tight', dpi=300, transparent=True)
        print(f"Plot saved to {save_path}")
        plt.close(fig)
    else:
        plt.show()

def save_kde_plot_image_with_cache(name, df, league_cache, output_dir=None):
    if output_dir is None:
        output_dir = os.path.join(os.getcwd(), "output", "catcher_plots")
    # Check if catcher has enough data before attempting to create plot
    filt_df = df[(df['PitcherTeam'] == team_code) & (df['Catcher'] == name)]
    if filt_df.empty:
        print(f"No data found for catcher {name}")
        return None
        
    proc_df = processing_full(filt_df)
    if proc_df.empty:
        print(f"No processed data for catcher {name}")
        return None
        
    # Check minimum sample sizes for KDE
    df_l = proc_df[proc_df['BatterSide'] == 'Left']
    df_r = proc_df[proc_df['BatterSide'] == 'Right']
    min_samples = 5
    
    print(f"DEBUG: {name} - LHB: {len(df_l)}, RHB: {len(df_r)}")
    
    if len(df_l) < min_samples or len(df_r) < min_samples:
        print(f"Catcher {name} has insufficient data for KDE plots (LHB: {len(df_l)}, RHB: {len(df_r)}). Skipping.")
        return None
    
    os.makedirs(output_dir, exist_ok=True)
    filepath = os.path.join(output_dir, f"{name.replace(',', '').replace(' ', '_')}.png")
    kde_plot_with_cache(name, df, league_cache, save_path=filepath)
    
    if os.path.exists(filepath):
        print(f"Image successfully saved: {filepath}")
        return filepath
    else:
        print(f"Failed to save image for catcher: {name}")
        return None

def PPTtoPDF_RemoveFirstSlide(inputFileName, outputFileName, formatType=32):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    # Ensure output ends with .pdf
    if not outputFileName.lower().endswith(".pdf"):
        outputFileName += ".pdf"

    # Open and modify the presentation
    deck = powerpoint.Presentations.Open(inputFileName, WithWindow=False)

    # Remove the first slide (assumes it's a template slide)
    if deck.Slides.Count >= 1:
        deck.Slides(1).Delete()

    # Save the modified presentation to PDF
    deck.SaveAs(outputFileName, formatType)  # 32 = PDF
    deck.Close()
    powerpoint.Quit()

    print(f"Saved PDF without template slide: {outputFileName}")

def export_each_catcher_slide_to_pdf(pptx_path, output_folder, catchers, date_str, template_slide_count=1):
    import time

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")

    presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)
    slide_count = presentation.Slides.Count

    os.makedirs(output_folder, exist_ok=True)

    # Start after the template slides (e.g., first slide is template)
    start_slide = template_slide_count + 1

    def safe_filename(name):
        return re.sub(r'[^A-Za-z0-9_\-]', '', name.replace(' ', '_'))

    try:
        for i, catcher in enumerate(catchers, start=start_slide):
            if i > slide_count:
                print(f"Slide index {i} exceeds total slides. Skipping.")
                continue

            # Create new presentation with the same template
            script_dir = os.path.dirname(os.path.abspath(__file__))
            template_path = os.path.join(script_dir, template_file)
            new_pres = powerpoint.Presentations.Open(template_path, WithWindow=False)
            
            # Copy slide from original
            presentation.Slides(i).Copy()

            # Paste slide into new presentation (replacing the template slide)
            time.sleep(0.2)  # slight delay to ensure clipboard is ready
            new_pres.Slides(1).Delete()  # Remove template slide
            new_pres.Slides.Paste()  # Paste the actual slide

            # Save new presentation as PDF
            pdf_filename = f"{safe_filename(catcher)}_{date_str}.pdf"
            pdf_path = os.path.abspath(os.path.join(output_folder, pdf_filename))

            new_pres.SaveAs(pdf_path, 32)  # 32 = PDF
            print(f"Saved PDF for {catcher} to {pdf_path}")

            new_pres.Close()
    except Exception as e:
        print(f"Error exporting slides to PDF: {e}")
    finally:
        presentation.Close()
        powerpoint.Quit()

def generate_pptx_from_catchers_template_with_cache(df, league_cache):
    safe_date = re.sub(r'[^A-Za-z0-9_]', '_', date)    
    base_folder = os.path.join(os.getcwd(), "output", f"catcher_report_{safe_date}")
    pdf_folder = os.path.join(base_folder, "pdfs")
    os.makedirs(pdf_folder, exist_ok=True)

    output_file = os.path.join(base_folder, "Catching_Report.pptx")

    df = df[df['PitcherTeam'] == team_code]
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_path = os.path.join(script_dir, template_file)

    prs = Presentation(template_path)
    slide_layout = prs.slide_layouts[0]

    catchers = df['Catcher'].dropna().unique()
    final_catchers = []  # track successfully processed catchers

    for catcher in catchers:
        # Check if catcher has enough data before processing
        filt = df[df['Catcher'] == catcher]
        calls = filt[filt['PitchCall'].isin(['BallCalled', 'StrikeCalled'])]
        calls_l = len(calls[calls['BatterSide'] == 'Left'])
        calls_r = len(calls[calls['BatterSide'] == 'Right'])
        
        if calls_l < 5 and calls_r < 5:
            print(f"Catcher {catcher} has insufficient data (LHB: {calls_l}, RHB: {calls_r}). Skipping.")
            continue
            
        path_to_img = save_kde_plot_image_with_cache(catcher, df, league_cache)
        if path_to_img is None or not os.path.exists(path_to_img):
            print(f"Image file for catcher {catcher} not found or insufficient data.")
            continue

        slide = prs.slides.add_slide(slide_layout)
        final_catchers.append(catcher)
        
        # Putting their name at the top
        name = catcher.split()
        full_name = name[-1] + ' ' + name[0]
        full_name = full_name[:-1]
        title = slide.shapes.title
        title.text = full_name
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.name = 'Beaver Bold'
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Adding a table for pop time/throw velo to 2b
        x, y, cx, cy = Inches(1.5), Inches(6.5), Inches(3.5), Inches(1)
        shape = slide.shapes.add_table(2, 2, x, y, cx, cy)
        table = shape.table
        table.columns[0].width = Inches(1.5)
        table.columns[1].width = Inches(2)
        
        # Finding the values
        filt = df[df['Catcher'] == catcher]
        throw2b = filt[filt['BasePositionX'] > 100]
        pop = throw2b['PopTime'].mean()
        arm = throw2b['ThrowSpeed'].mean()

        # NCAA average pop time and arm speed on throws to 2b
        lgpop = 2.1288857946774464
        lgarm = 77.27192782800913

        # Adding values to table
        table.cell(0,0).text = 'Pop Time'
        table.cell(1,0).text = 'Throw Velo'        

        # Finding diff from league average
        pop_diff = (pop - lgpop) / lgpop
        arm_diff = (arm - lgarm) / lgarm

        # Colormapping
        cmap = colors.LinearSegmentedColormap.from_list("bwr", ["blue", "white", "red"])
        norm = colors.Normalize(vmin=-0.1, vmax=0.1)
        def get_rgb_color(diff, invert=False):
            """Returns an RGBColor mapped from a normalized diff. 
            If invert=True, better performance is lower (e.g., Pop Time)."""
            if invert:
                diff = -diff
            rgba = cmap(norm(diff))
            r, g, b = [int(255 * c) for c in rgba[:3]]
            return RGBColor(r, g, b)
        
        # Set value and background color for Pop Time
        cell_pop = table.cell(0,1)
        cell_pop.text = f"{pop:.2f} ({-pop_diff:+.1%})"
        cell_pop.fill.solid()
        cell_pop.fill.fore_color.rgb = get_rgb_color(pop_diff, invert=True)  # lower is better

        # Set value and background color for Throw Velo
        cell_arm = table.cell(1,1)
        cell_arm.text = f"{arm:.2f} ({arm_diff:+.1%})"
        cell_arm.fill.solid()
        cell_arm.fill.fore_color.rgb = get_rgb_color(arm_diff, invert=False)  # higher is better

        # Apply bold black text only (do NOT reset fill color)
        for row in table.rows:
            cell = row.cells[0]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

        # Add SB/Attempt textbox
        left, top, width, height = Inches(6.1), Inches(6.75), Inches(3.4), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        # Set grey background fill (e.g., RGB(200, 200, 200))
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(200, 200, 200)
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        sb, att = sb_att_dict.get(catcher, (0, 0))  # Default to (0, 0) if catcher not in dict
        rate = sb/att * 100 if att !=0 else 0
        p.text = f"SB/Attempt: {sb}/{att} ({rate:.1f}%)"
        p.alignment = PP_ALIGN.CENTER

        # Apply font styling
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)  # Black

        # add date section
        left, top, width, height = Inches(8), Inches(0), Inches(3), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = date
        p1.alignment = PP_ALIGN.RIGHT

        # Apply font styling
        run1 = p1.runs[0]
        run1.font.size = Pt(18)
        run1.font.name = 'Beaver Bold'
        run1.font.color.rgb = RGBColor(0, 0, 0)  # Black

        # Find number of pitches to each handedness
        calls = filt[filt['PitchCall'].isin(['BallCalled', 'StrikeCalled'])]
        calls_l = len(calls[calls['BatterSide'] == 'Left'])
        calls_r = len(calls[calls['BatterSide'] == 'Right'])

        # Adding text for each handedness
        left, top, width, height = Inches(1), Inches(1.1), Inches(4), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        l1 = text_frame.paragraphs[0]
        l1.text = f"LHB: {calls_l} pitches"
        l1.alignment = PP_ALIGN.LEFT
        runl1 = l1.runs[0]
        runl1.font.size = Pt(18)
        runl1.font.name = 'Beaver Bold'
        runl1.font.color.rgb = RGBColor(0, 0, 0)  # Black

        left, top, width, height = Inches(6), Inches(1.1), Inches(4), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        r1 = text_frame.paragraphs[0]
        r1.text = f"RHB: {calls_r} pitches"
        r1.alignment = PP_ALIGN.RIGHT
        runr1 = r1.runs[0]
        runr1.font.size = Pt(18)
        runr1.font.name = 'Beaver Bold'
        runr1.font.color.rgb = RGBColor(0, 0, 0)  # Black

        for shape in slide.shapes:
            if not shape.is_placeholder:
                continue

            placeholder_format = shape.placeholder_format

            # If the placeholder is of type PICTURE, insert the image
            if placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                # Save original position and size
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Remove the placeholder
                sp = shape._element
                slide.shapes._spTree.remove(sp)

                # Insert picture in its place
                slide.shapes.add_picture(path_to_img, left, top, width=width, height=height)

    prs.save(output_file)
    output_file = os.path.abspath(output_file)
    print(f"PowerPoint saved to: {output_file}")

    pdf_output = os.path.join(base_folder, "Catching_Report.pdf")
    PPTtoPDF_RemoveFirstSlide(output_file, pdf_output)

    # Save individual PDFs into pdfs/ subfolder
    pdf_folder_individual = pdf_folder
    os.makedirs(pdf_folder_individual, exist_ok=True)

    catchers = df['Catcher'].dropna().unique()
    export_each_catcher_slide_to_pdf(output_file, pdf_folder_individual, catchers, safe_date)

# Load the pre-computed KDE cache
cache_file = "kde_cache.npz"
league_cache = load_kde_cache(cache_file)

if league_cache is None:
    print("Exiting due to missing cache file.")
    exit(1)

print("✅ KDE cache loaded successfully!")
print("Generating catcher reports...")

# Generate reports using cached league data
generate_pptx_from_catchers_template_with_cache(global_df, league_cache) 